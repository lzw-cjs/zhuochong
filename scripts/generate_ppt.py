"""使用 SenseNova API 生成项目介绍 PPT"""
import os
import sys
import json
import requests
from pathlib import Path
from dotenv import load_dotenv

# 加载 .env
load_dotenv()

BASE_URL = os.getenv("SNSDK_BASE_URL", "https://token.sensenova.cn/v1")
API_KEY = os.getenv("SNSDK_API_KEY", "")

if not API_KEY:
    print("错误：请在 .env 文件中设置 SNSDK_API_KEY")
    sys.exit(1)


def call_llm(prompt: str, system_prompt: str = "") -> str:
    """调用 SenseNova LLM API"""
    url = f"{BASE_URL}/chat/completions"
    headers = {
        "Authorization": f"Bearer {API_KEY}",
        "Content-Type": "application/json",
    }

    messages = []
    if system_prompt:
        messages.append({"role": "system", "content": system_prompt})
    messages.append({"role": "user", "content": prompt})

    payload = {
        "model": "sensenova-6.7-flash-lite",
        "messages": messages,
        "temperature": 0.7,
        "max_tokens": 4096,
    }

    resp = requests.post(url, headers=headers, json=payload, timeout=60)
    if resp.status_code != 200:
        print(f"LLM API 错误: {resp.status_code}")
        print(resp.text)
        return ""

    data = resp.json()
    return data["choices"][0]["message"]["content"]


def generate_ppt_content(project_name: str, pages: int = 10) -> list[dict]:
    """用 LLM 生成 PPT 各页内容"""
    system_prompt = """你是一个专业的 PPT 内容策划师。
请根据用户提供的项目信息，生成 PPT 的每一页内容。

输出格式为 JSON 数组，每个元素包含：
- title: 页面标题
- bullets: 要点列表（3-5 个）
- notes: 演讲备注

确保内容专业、简洁、有逻辑性。"""

    user_prompt = f"""请为"{project_name}"项目生成一份 {pages} 页的介绍 PPT 内容。

项目信息：
- 技术栈：Python + PySide6 + JSON Storage
- 核心功能：像素风桌面宠物、动画系统、聊天互动、日程管理、提醒系统
- 架构：分层架构（UI层、交互层、逻辑层、数据层、基础设施层）

页面结构建议：
1. 封面
2. 项目概述
3. 核心功能
4. 技术架构
5. 动画系统
6. 聊天系统
7. 日程管理
8. 提醒系统
9. 演示效果
10. 总结与展望

请输出 JSON 格式。"""

    print("正在用 LLM 生成 PPT 内容...")
    result = call_llm(user_prompt, system_prompt)

    # 尝试解析 JSON
    try:
        # 提取 JSON 部分
        if "```json" in result:
            result = result.split("```json")[1].split("```")[0]
        elif "```" in result:
            result = result.split("```")[1].split("```")[0]
        return json.loads(result.strip())
    except json.JSONDecodeError:
        print("JSON 解析失败，使用默认内容")
        return None


def create_ppt(slides_content: list[dict], output_path: str, project_name: str):
    """使用 python-pptx 创建 PPT"""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 颜色主题
    PRIMARY = RGBColor(0x1A, 0x73, 0xE8)    # 蓝色
    DARK = RGBColor(0x20, 0x2A, 0x44)       # 深蓝
    LIGHT_BG = RGBColor(0xF8, 0xF9, 0xFA)   # 浅灰背景
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)

    for i, slide_data in enumerate(slides_content):
        if i == 0:
            # 封面页
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
            bg = slide.background
            fill = bg.fill
            fill.solid()
            fill.fore_color.rgb = PRIMARY

            # 标题
            title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(2))
            title_frame = title_box.text_frame
            title_frame.text = slide_data.get("title", project_name)
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(44)
            title_para.font.color.rgb = WHITE
            title_para.font.bold = True
            title_para.alignment = PP_ALIGN.CENTER

            # 副标题
            sub_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11), Inches(1))
            sub_frame = sub_box.text_frame
            sub_frame.text = "智能桌面宠物应用 | Python + PySide6"
            sub_para = sub_frame.paragraphs[0]
            sub_para.font.size = Pt(20)
            sub_para.font.color.rgb = RGBColor(0xBB, 0xDE, 0xFB)
            sub_para.alignment = PP_ALIGN.CENTER
        else:
            # 内容页
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            # 顶部色条
            shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.08))
            shape.fill.solid()
            shape.fill.fore_color.rgb = PRIMARY
            shape.line.fill.background()

            # 标题
            title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(11), Inches(0.8))
            title_frame = title_box.text_frame
            title_frame.text = slide_data.get("title", f"第 {i+1} 页")
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(32)
            title_para.font.color.rgb = DARK
            title_para.font.bold = True

            # 要点内容
            bullets = slide_data.get("bullets", [])
            if bullets:
                content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11), Inches(5))
                content_frame = content_box.text_frame
                content_frame.word_wrap = True

                for j, bullet in enumerate(bullets):
                    if j == 0:
                        para = content_frame.paragraphs[0]
                    else:
                        para = content_frame.add_paragraph()
                    para.text = f"• {bullet}"
                    para.font.size = Pt(18)
                    para.font.color.rgb = DARK
                    para.space_after = Pt(12)
                    para.space_before = Pt(6)

            # 备注
            notes = slide_data.get("notes", "")
            if notes:
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = notes

        print(f"  已创建第 {i+1} 页: {slide_data.get('title', '')}")

    prs.save(output_path)
    print(f"\nPPT 已保存: {output_path}")


# 默认 PPT 内容（LLM 调用失败时使用）
DEFAULT_SLIDES = [
    {
        "title": "智能桌面宠物 Smart Desktop Pet",
        "bullets": ["像素风桌面宠物应用", "Python + PySide6 技术栈", "集成日程管理与提醒功能"],
        "notes": "封面页"
    },
    {
        "title": "项目概述",
        "bullets": [
            "一款运行在桌面的像素风宠物应用",
            "宠物可互动：拖拽、点击、右键菜单",
            "内置聊天系统，支持关键词回复",
            "集成日程管理，支持事件创建和提醒",
            "像素风动画，5 种状态：待机、行走、睡眠、开心、警报"
        ],
        "notes": "介绍项目的核心价值和目标用户"
    },
    {
        "title": "核心功能",
        "bullets": [
            "动画系统：5 种状态，帧动画控制，状态机管理",
            "交互系统：拖拽移动、右键菜单、气泡提示",
            "聊天引擎：基于规则的关键词匹配回复",
            "日程管理：事件 CRUD、多日历支持、月历视图",
            "提醒系统：定时轮询、音效播放、超时检测"
        ],
        "notes": "详细介绍五大核心功能模块"
    },
    {
        "title": "技术架构",
        "bullets": [
            "UI 表现层：Window、Tray、Bubble、ChatPanel、SchedulePanel",
            "交互控制层：Behavior、ChatEngine、Drag、ContextMenu",
            "核心逻辑层：States、Animator、ReminderEngine、OverdueDetector",
            "数据存储层：Settings、Event、Calendar、JsonStore",
            "基础设施层：SoundManager、Assets、Migration"
        ],
        "notes": "分层架构设计，各层职责清晰"
    },
    {
        "title": "动画系统",
        "bullets": [
            "SpriteAnimator：帧动画控制器，QTimer 驱动",
            "PetState 状态机：IDLE/WALK/SLEEP/HAPPY/ALERT",
            "状态转换规则：定义合法的状态切换路径",
            "占位帧生成：32x32 像素风树懒图",
            "帧间隔：不同状态不同速度（150ms-800ms）"
        ],
        "notes": "动画系统的技术实现细节"
    },
    {
        "title": "聊天系统",
        "bullets": [
            "ChatEngine 抽象基类，支持策略模式替换",
            "RuleBasedEngine：基于 chat_rules.json 的关键词匹配",
            "架构预留：可替换为 LLM 引擎",
            "UI：QScrollArea + QLineEdit + QPushButton",
            "支持多轮对话和历史记录"
        ],
        "notes": "聊天系统的设计和扩展性"
    },
    {
        "title": "日程管理",
        "bullets": [
            "Event 数据模型：标题、时间、分类、优先级、截止时间",
            "ScheduleStore：JSON 文件存储，支持导入导出",
            "CalendarGrid：月历网格组件，事件日期高亮",
            "EventDialog：事件编辑对话框，支持提醒设置",
            "多格式导入：Markdown/纯文本/JSON"
        ],
        "notes": "日程管理的完整功能链"
    },
    {
        "title": "提醒系统",
        "bullets": [
            "ReminderEngine：60 秒轮询，检测待触发事件",
            "去重机制：避免重复提醒",
            "抑制功能：临时关闭提醒",
            "SoundManager：QSoundEffect 音效播放",
            "ALERT 状态：抖动动画 + 红色感叹号"
        ],
        "notes": "提醒系统的工作流程"
    },
    {
        "title": "演示效果",
        "bullets": [
            "宠物在桌面自由移动，可拖拽到任意位置",
            "点击宠物触发互动，右键弹出功能菜单",
            "聊天窗口支持实时对话",
            "日程面板显示月历和事件列表",
            "提醒触发时宠物变为警报状态并播放音效"
        ],
        "notes": "实际运行效果展示"
    },
    {
        "title": "总结与展望",
        "bullets": [
            "v1.0 已完成：6 个阶段、25 个计划全部交付",
            "技术亮点：分层架构、状态机、原子写入、Schema 迁移",
            "未来方向：LLM 聊天引擎、更多动画状态、云端同步",
            "开源地址：github.com/OpenSenseNova/SenseNova-Skills",
            "感谢观看！"
        ],
        "notes": "总结项目成果和未来规划"
    }
]


if __name__ == "__main__":
    project_name = "智能桌面宠物 Smart Desktop Pet"
    output_dir = Path("output")
    output_dir.mkdir(exist_ok=True)
    output_path = output_dir / "project_intro.pptx"

    # 尝试用 LLM 生成内容
    slides = generate_ppt_content(project_name, pages=10)

    if not slides:
        print("使用默认 PPT 内容")
        slides = DEFAULT_SLIDES

    # 生成 PPT
    print(f"\n正在创建 PPT（{len(slides)} 页）...")
    create_ppt(slides, str(output_path), project_name)
